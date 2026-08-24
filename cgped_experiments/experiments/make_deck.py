# -*- coding: utf-8 -*-
"""Builds CG-PED_Defense_Deck.pptx — the narrative defense deck."""
import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR, MSO_AUTO_SIZE
from pptx.enum.shapes import MSO_SHAPE

HERE = os.path.dirname(os.path.abspath(__file__))
ROOT = os.path.abspath(os.path.join(HERE, os.pardir, os.pardir))
FIG  = os.path.join(ROOT, "cgped_experiments", "figures")
SCR  = FIG
OUT  = os.path.join(ROOT, "CG-PED_Defense_Deck.pptx")

# ----------------------------------------------------------------- palette
INK    = "16202E"   # near-navy ink on light slides
BODY   = "334155"
MUTED  = "708196"
PAPER  = "FFFFFF"
SOFT   = "F4F6FA"
LINE   = "DCE3EC"
DARK   = "0F1B2D"   # dark statement slides
DARKSF = "1B2A42"
SNOW   = "F1F5F9"

AMBER  = "E08A2E"   # curriculum
RED    = "C0392B"   # LoRA
TEAL   = "0E8074"   # distillation
BLUE   = "2563EB"   # baseline
GOLD   = "F0B429"

FONT = "Arial"
MONO = "Courier New"

W, H = Inches(13.333), Inches(7.5)

prs = Presentation()
prs.slide_width, prs.slide_height = W, H
BLANK = prs.slide_layouts[6]

_n = [0]


# ----------------------------------------------------------------- helpers
def rgb(h):
    return RGBColor.from_string(h)


def slide(dark=False, soft=False):
    s = prs.slides.add_slide(BLANK)
    col = DARK if dark else (SOFT if soft else PAPER)
    r = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, W, H)
    r.fill.solid(); r.fill.fore_color.rgb = rgb(col)
    r.line.fill.background(); r.shadow.inherit = False
    s._dark = dark
    return s


def box(s, x, y, w, h, color, line=None, lw=1.25, shape=MSO_SHAPE.RECTANGLE,
        radius=None):
    sh = s.shapes.add_shape(shape, Inches(x), Inches(y), Inches(w), Inches(h))
    if color is None:
        sh.fill.background()
    else:
        sh.fill.solid(); sh.fill.fore_color.rgb = rgb(color)
    if line:
        sh.line.color.rgb = rgb(line); sh.line.width = Pt(lw)
    else:
        sh.line.fill.background()
    sh.shadow.inherit = False
    if radius is not None and shape == MSO_SHAPE.ROUNDED_RECTANGLE:
        try:
            sh.adjustments[0] = radius
        except Exception:
            pass
    return sh


def text(s, x, y, w, h, runs, size=18, color=None, bold=False, align="l",
         font=FONT, space=1.0, anchor="t", space_after=0):
    """runs: str, or list of (text, dict-of-overrides) tuples, or list of
    paragraph-lists."""
    tb = s.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.auto_size = MSO_AUTO_SIZE.NONE
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    tf.vertical_anchor = {"t": MSO_ANCHOR.TOP, "m": MSO_ANCHOR.MIDDLE,
                          "b": MSO_ANCHOR.BOTTOM}[anchor]
    default_col = color or (SNOW if getattr(s, "_dark", False) else INK)

    if isinstance(runs, str):
        paras = [[(runs, {})]]
    elif runs and isinstance(runs[0], tuple):
        paras = [runs]
    else:
        paras = runs

    for i, para in enumerate(paras):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = {"l": PP_ALIGN.LEFT, "c": PP_ALIGN.CENTER,
                       "r": PP_ALIGN.RIGHT, "j": PP_ALIGN.JUSTIFY}[align]
        p.line_spacing = space
        if space_after:
            p.space_after = Pt(space_after)
        if isinstance(para, str):
            para = [(para, {})]
        for t, o in para:
            r = p.add_run(); r.text = t
            f = r.font
            f.name = o.get("font", font)
            f.size = Pt(o.get("size", size))
            f.bold = o.get("bold", bold)
            f.italic = o.get("italic", False)
            f.color.rgb = rgb(o.get("color", default_col))
    return tb


def kicker(s, t, color=None, x=0.85, y=0.62):
    text(s, x, y, 11, 0.3, t.upper(), size=12.5, bold=True,
         color=color or (GOLD if getattr(s, "_dark", False) else MUTED))


def title(s, t, y=1.02, size=38, color=None, x=0.85, w=11.6, h=1.2):
    text(s, x, y, w, h, t, size=size, bold=True,
         color=color or (PAPER if getattr(s, "_dark", False) else INK),
         space=1.03)


def rule(s, x=0.85, y=None, w=1.4, color=GOLD, h=0.055):
    box(s, x, y, w, h, color)


def bullets(s, x, y, w, items, size=17.5, gap=0.44, color=None, dot=None,
            dotsize=0.105, lh=1.18, wrap_h=0.42):
    """items: list of str or (str, extra_runs_list)."""
    col = color or (SNOW if getattr(s, "_dark", False) else BODY)
    dcol = dot or (GOLD if getattr(s, "_dark", False) else INK)
    yy = y
    for it in items:
        box(s, x, yy + 0.115, dotsize, dotsize, dcol, shape=MSO_SHAPE.OVAL)
        if isinstance(it, str):
            it = [(it, {})]
        text(s, x + 0.32, yy, w - 0.32, wrap_h, it, size=size, color=col,
             space=lh)
        yy += gap
    return yy


def footer(s, label=""):
    _n[0] += 1
    d = getattr(s, "_dark", False)
    text(s, 0.85, 6.92, 8.0, 0.3, label, size=10.5,
         color="53627A" if d else "9AA7B8")
    text(s, 11.0, 6.92, 1.5, 0.3, f"{_n[0]:02d}", size=10.5, align="r",
         color="53627A" if d else "9AA7B8")


def pic(s, path, x, y, w):
    im = s.shapes.add_picture(path, Inches(x), Inches(y), width=Inches(w))
    return im


def card(s, x, y, w, h, accent, head, sub, body, headsize=21, bodysize=14.5,
         fill=PAPER, edge=LINE):
    box(s, x, y, w, h, fill, line=edge)
    box(s, x, y, w, 0.085, accent)
    text(s, x + 0.34, y + 0.42, w - 0.68, 0.42, head, size=headsize, bold=True,
         color=accent)
    if sub:
        text(s, x + 0.34, y + 0.95, w - 0.68, 0.4, sub, size=13,
             color=MUTED, bold=True)
    text(s, x + 0.34, y + 1.42, w - 0.68, h - 1.8, body, size=bodysize,
         color=BODY, space=1.22)


def chip(s, x, y, w, h, label, fill, txtcol=PAPER, size=13.5, bold=True,
         radius=0.35):
    sh = box(s, x, y, w, h, fill, shape=MSO_SHAPE.ROUNDED_RECTANGLE,
             radius=radius)
    text(s, x, y, w, h, label, size=size, bold=bold, color=txtcol,
         align="c", anchor="m")
    return sh


def statnum(s, x, y, w, num, label, color, numsize=44, labsize=12.5,
            gap=0.62):
    text(s, x, y, w, gap, num, size=numsize, bold=True, color=color,
         align="c", space=0.95)
    text(s, x, y + gap, w, 0.72, label, size=labsize, color=MUTED, align="c",
         space=1.15)


# =========================================================== 01 TITLE
s = slide(dark=True)
box(s, 0, 0, 0.28, 7.5, GOLD)
text(s, 1.0, 1.28, 11, 0.34, "RESEARCH DEFENSE  ·  ABEL TEZARE", size=13,
     bold=True, color=GOLD)
text(s, 1.0, 1.95, 11.2, 2.2,
     [[("Do three cost-saving\ntricks ", {}), ("add up", {"color": GOLD}),
       ("?", {})]], size=52, bold=True, color=PAPER, space=1.06)
box(s, 1.0, 4.22, 1.4, 0.055, GOLD)
text(s, 1.0, 4.58, 10.6, 1.0,
     [[("CG-PED", {"bold": True, "color": PAPER}),
       ("  —  Curriculum Learning  ×  LoRA  ×  Knowledge Distillation, "
        "tested as a full 2³ factorial", {})]],
     size=19, color="B9C6D8", space=1.25)
for i, (num, lab) in enumerate([("48", "training runs"), ("8", "factorial cells"),
                                ("3", "seeds"), ("2", "benchmarks"),
                                ("1", "CPU core")]):
    x = 1.0 + i * 2.16
    text(s, x, 5.65, 1.9, 0.55, num, size=30, bold=True, color=PAPER)
    text(s, x, 6.18, 1.9, 0.3, lab, size=11.5, color="7D8FA8")
footer(s, "CG-PED")

# =========================================================== 02 CONFESSION
s = slide()
kicker(s, "Before the science")
title(s, "Where this project actually came from")
rule(s, y=2.05)
bullets(s, 0.85, 2.55, 7.3, [
    [("Six months ago my world was ", {}),
     ("Python, DSA and web development", {"bold": True, "color": INK}),
     (". Machine learning was a word I had heard.", {})],
    [("Everything in this deck — the design, the code, the bugs, the "
      "negative result — came out of ", {}),
     ("the training and the projects you gave us", {"bold": True, "color": INK}),
     (".", {})],
    [("I started behind the other trainees. I caught up. And I did not stop "
      "at the assignment.", {})],
], size=17, gap=1.05, wrap_h=0.95)
box(s, 8.55, 2.35, 3.9, 2.6, DARK)
text(s, 8.95, 2.75, 3.1, 1.9,
     [[("“This deck is\nthe receipt.”", {})]],
     size=25, bold=True, color=PAPER, space=1.15)
text(s, 8.95, 4.36, 3.1, 0.4, "— and at the end, what I want to build next",
     size=11.5, color="8FA0B6")
footer(s, "Introduction")

# =========================================================== 03 THE MYSTERY
s = slide(dark=True)
kicker(s, "Motivation")
text(s, 0.85, 1.62, 11.6, 2.8,
     [[("In 2019 a decent model needed\na data centre.", {"color": "8CA0BA"})],
      [("In 2025 people fine-tune one\non a laptop.", {"color": PAPER})]],
     size=36, bold=True, space=1.14, space_after=12)
box(s, 0.85, 4.92, 1.4, 0.055, GOLD)
text(s, 0.85, 5.35, 11.2, 1.2,
     [[("Nothing magic happened in between. ", {"color": "B9C6D8"}),
       ("A pile of efficiency techniques happened", {"bold": True, "color": GOLD}),
       (" — and almost everyone reports them one at a time.", {"color": "B9C6D8"})]],
     size=19.5, space=1.3)
footer(s, "Motivation")

# =========================================================== 04 WHY MIDDLE
s = slide()
kicker(s, "Motivation")
title(s, "Why I deliberately studied “old” techniques")
# timeline
ty = 2.35
box(s, 0.95, ty + 0.30, 11.4, 0.035, LINE)
stops = [("2009", "Curriculum\nBengio et al.", AMBER),
         ("2015", "Distillation\nHinton et al.", TEAL),
         ("2021", "LoRA\nHu et al.", RED),
         ("2024+", "today's efficient\nAI stack", MUTED)]
for i, (yr, nm, c) in enumerate(stops):
    x = 1.05 + i * 3.32
    box(s, x, ty + 0.19, 0.26, 0.26, c, shape=MSO_SHAPE.OVAL)
    text(s, x + 0.42, ty + 0.12, 2.6, 0.4, yr, size=17, bold=True, color=c)
    text(s, x + 0.42, ty + 0.55, 2.7, 0.75, nm, size=13, color=MUTED, space=1.2)
bullets(s, 0.85, 3.85, 11.5, [
    [("Recent tricks are still unstable. If a combination behaved strangely, "
      "I could not tell whether that was ", {}),
     ("the interaction or the technique", {"bold": True, "color": INK}), (".", {})],
    [("These three have well-characterised solo effects. So any surprise is "
      "a ", {}), ("finding", {"bold": True, "color": INK}),
     (", not a mystery.", {})],
    [("They are the load-bearing beams under today's efficient AI — and "
      "practitioners under a budget still stack all three at once.", {})],
], size=16.5, gap=0.74, wrap_h=0.68)
box(s, 0.85, 6.05, 11.5, 0.62, SOFT)
text(s, 1.15, 6.05, 11.0, 0.62,
     [[("The contribution is the design. The techniques are just the test "
        "subjects.", {})]], size=15.5, bold=True, color=INK, anchor="m")
footer(s, "Motivation")

# =========================================================== 05 THE CAST
s = slide(soft=True)
kicker(s, "The cast")
title(s, "Three different ways to spend less")
cards = [
    (AMBER, "C  ·  Curriculum", "Teach easy before hard",
     "You don't open a maths course with calculus.\n\nSort the data by "
     "difficulty, start the model on 15% of it, open the pool as it copes.\n\n"
     "Price: free. Just reorder."),
    (RED, "L  ·  LoRA", "Sticky notes, not a rewrite",
     "Freeze the book. Write small rank-4 notes in the margin.\n\n"
     "Only the notes get trained — 4× fewer weights to update.\n\n"
     "Price: memory saved. Capacity spent."),
    (TEAL, "D  ·  Distillation", "Learn from a senior",
     "A bigger teacher already made the mistakes. The student copies its "
     "soft answers, not just the labels.\n\n"
     "Price: you have to pay for the teacher too."),
]
for i, (c, h, sub, b) in enumerate(cards):
    card(s, 0.85 + i * 3.92, 2.28, 3.62, 4.05, c, h, sub, b)
footer(s, "Background")

# =========================================================== 06 THE COUPONS
s = slide(dark=True)
kicker(s, "The untested assumption")
text(s, 0.85, 1.55, 7.3, 2.4,
     [[("You reach the checkout holding a coupon, a loyalty card, and a "
        "cashback offer.", {"color": "B9C6D8"})],
      [("Everyone assumes 10% + 15% + 5% = 30% off.", {"color": PAPER})]],
     size=27, bold=True, space=1.2, space_after=14)
box(s, 0.85, 4.82, 1.4, 0.055, GOLD)
text(s, 0.85, 5.22, 7.4, 1.9,
     [[("Nobody ran\nthe receipt.", {})]], size=42, bold=True, color=GOLD,
     space=1.08)
box(s, 8.6, 1.55, 3.88, 5.05, DARKSF)
text(s, 8.95, 1.95, 3.2, 0.4, "WHAT THE LITERATURE DOES", size=12, bold=True,
     color=GOLD)
bullets(s, 8.95, 2.55, 3.2, [
    "CLPD (2026): curriculum + distillation",
    "POCL (2025): curriculum + distillation",
    "TSCL (2026): sequenced KD objectives",
], size=13.5, gap=0.62, color="B9C6D8", dot=GOLD, dotsize=0.085, wrap_h=0.6)
box(s, 8.95, 4.42, 3.2, 0.02, "3B4C68")
text(s, 8.95, 4.68, 3.2, 1.7,
     [[("Every one of them compares the combination to ", {"color": "8CA0BA"}),
       ("nothing", {"bold": True, "color": PAPER}),
       (".\n\nThat answers “better than nothing?” — not "
        "“better than its parts?”", {"color": "8CA0BA"})]],
     size=13.5, space=1.25)
footer(s, "The gap")

# =========================================================== 07 THE MATH
s = slide()
kicker(s, "The idea in one line")
title(s, "Composition is not combination")
box(s, 0.85, 2.05, 11.5, 1.28, DARK)
text(s, 0.85, 2.05, 11.5, 1.28,
     [[("f(g(h(x + y + z)))", {"color": PAPER}),
       ("   ≠   ", {"color": GOLD}),
       ("f(x) + g(y) + h(z)", {"color": PAPER})]],
     size=32, bold=True, align="c", anchor="m", font=MONO)
box(s, 0.85, 3.62, 5.55, 1.62, SOFT)
text(s, 1.15, 3.85, 4.95, 1.2,
     [[("COMBINATION", {"bold": True, "color": TEAL, "size": 13})],
      [("Effects sit side by side. You can add them "
        "and the arithmetic is honest.", {"color": BODY, "size": 15})]],
     size=15, space=1.25, space_after=7)
box(s, 6.8, 3.62, 5.55, 1.62, SOFT)
text(s, 7.1, 3.85, 4.95, 1.2,
     [[("COMPOSITION", {"bold": True, "color": RED, "size": 13})],
      [("Each technique changes the input of the next. "
        "The sum is a guess.", {"color": BODY, "size": 15})]],
     size=15, space=1.25, space_after=7)
bullets(s, 0.85, 5.48, 11.5, [
    [("The curriculum changes ", {}), ("what LoRA sees", {"bold": True, "color": INK}),
     (". LoRA changes ", {}),
     ("how much the student can absorb", {"bold": True, "color": INK}),
     (" from the teacher.", {})],
    [("Training with all three is composition. Reporting them as a sum "
      "assumes combination. ", {}),
     ("That gap is the whole paper.", {"bold": True, "color": RED})],
], size=15.5, gap=0.72, wrap_h=0.68)
footer(s, "The research question")

# =========================================================== 08 THREE WORLDS
s = slide(soft=True)
kicker(s, "Why one-factor-at-a-time cannot answer it")
title(s, "Three worlds are possible. Four experiments see one.")
pic(s, os.path.join(SCR, "fig_worlds.png"), 0.85, 2.42, 11.5)
box(s, 0.85, 5.62, 11.5, 1.0, DARK)
text(s, 1.25, 5.62, 10.7, 1.0,
     [[("Run C alone, L alone, D alone and you get three bars — and then you "
        "assume the fourth. ", {"color": "B9C6D8"}),
       ("You have to actually run the combination.", {"bold": True, "color": PAPER})]],
     size=16, anchor="m", space=1.25)
footer(s, "The research question")

# =========================================================== 09 THE DESIGN
s = slide()
kicker(s, "Design")
title(s, "A full 2³ factorial — all eight cells", size=36)
cells = [("- - -", "naive baseline", [0, 0, 0]), ("C - -", "curriculum", [1, 0, 0]),
         ("- L -", "LoRA", [0, 1, 0]), ("- - D", "distillation", [0, 0, 1]),
         ("C L -", "curr + LoRA", [1, 1, 0]), ("C - D", "curr + distil", [1, 0, 1]),
         ("- L D", "LoRA + distil", [0, 1, 1]), ("C L D", "full CG-PED", [1, 1, 1])]
cw, ch, gx, gy = 2.72, 1.32, 2.92, 1.52
for i, (nm, lab, on) in enumerate(cells):
    x = 0.85 + (i % 4) * gx
    y = 2.15 + (i // 4) * gy
    is_naive, is_full = (i == 0), (i == 7)
    fill = DARK if is_full else (PAPER if not is_naive else PAPER)
    edge = RED if is_full else (BLUE if is_naive else LINE)
    box(s, x, y, cw, ch, fill, line=edge, lw=2.0 if (is_full or is_naive) else 1.25)
    text(s, x + 0.25, y + 0.22, cw - 0.5, 0.42, nm, size=21, bold=True,
         font=MONO, color=PAPER if is_full else (BLUE if is_naive else INK))
    text(s, x + 0.25, y + 0.70, cw - 0.5, 0.32, lab, size=12,
         color="9FB0C6" if is_full else MUTED)
    for j, c in enumerate([AMBER, RED, TEAL]):
        box(s, x + cw - 0.95 + j * 0.28, y + 0.30, 0.19, 0.19,
            c if on[j] else ("34455F" if is_full else "E2E8F0"),
            shape=MSO_SHAPE.OVAL)
box(s, 0.85, 5.34, 11.5, 0.03, LINE)
eq = [("8", "cells"), ("×", ""), ("3", "seeds"), ("×", ""),
      ("2", "benchmarks"), ("=", ""), ("48", "training runs")]
xx = 0.95
for num, lab in eq:
    wdt = 1.15 if lab else 0.5
    text(s, xx, 5.6, wdt, 0.6, num, size=31 if lab else 24, bold=True,
         color=RED if num == "48" else INK, align="c")
    if lab:
        text(s, xx - 0.4, 6.16, wdt + 0.8, 0.32, lab, size=11.5, color=MUTED,
             align="c")
    xx += wdt + 0.25
text(s, 9.35, 5.6, 3.0, 1.1,
     [[("Everything else held identical: architecture, learning rate, batch "
        "size, step budget, seeds.", {})]], size=13, color=MUTED, space=1.25)
footer(s, "Design")

# =========================================================== 10 WHY 8 NOT 4
s = slide(soft=True)
kicker(s, "Design")
title(s, "Why eight cells, when four would look cheaper")
box(s, 0.85, 2.3, 5.55, 3.5, PAPER, line=LINE)
text(s, 1.2, 2.65, 4.9, 0.45, "FOUR CELLS", size=15, bold=True, color=MUTED)
text(s, 1.2, 3.15, 4.9, 0.5, "one-factor-at-a-time", size=20, bold=True,
     color=INK)
bullets(s, 1.2, 3.95, 4.9, [
    "each factor's effect vs baseline",
    "no interaction term",
    "no independent noise estimate",
], size=14.5, gap=0.52, dot=MUTED, dotsize=0.09, wrap_h=0.5)
box(s, 6.8, 2.3, 5.55, 3.5, DARK)
text(s, 7.15, 2.65, 4.9, 0.45, "EIGHT CELLS", size=15, bold=True, color=GOLD)
text(s, 7.15, 3.15, 4.9, 0.5, "full factorial", size=20, bold=True, color=PAPER)
bullets(s, 7.15, 3.85, 4.9, [
    "every main effect, averaged over 4 pairs",
    "the additive prediction for C+L+D",
    "observed − predicted = the composition gap",
    "seed spread gives the noise floor",
], size=14.5, gap=0.5, color="B9C6D8", dot=GOLD, dotsize=0.09, wrap_h=0.48)
text(s, 0.85, 6.15, 11.5, 0.7,
     [[("Exactly double the cost. It buys the one quantity no cheaper design "
        "can produce.", {})]], size=17, bold=True, color=INK, align="c")
footer(s, "Design")

# =========================================================== 11 THE BENCH
s = slide()
kicker(s, "Setup")
title(s, "The bench — small enough to run 48 times", size=36)
rows = [
    ("Student", "2-layer Transformer, d_model 64 — 82K parameters, "
     "character-level", INK),
    ("Teacher", "4-layer, 561K parameters — trained from scratch on the same "
     "data, then frozen. Not a pretrained model.", INK),
    ("Budget", "300 steps, LR 4e-3, batch 32, evaluated every 15 steps", INK),
    ("Benchmark A", "Synthetic SC-LM — difficulty is known by construction, "
     "noise-free", AMBER),
    ("Benchmark B", "Tiny Shakespeare — difficulty comes from the corpus, "
     "not from me", AMBER),
]
yy = 2.2
for lab, val, c in rows:
    text(s, 0.85, yy, 2.1, 0.4, lab, size=14.5, bold=True, color=c)
    text(s, 3.1, yy, 6.1, 0.62, val, size=15, color=BODY, space=1.2)
    box(s, 0.85, yy + 0.66, 8.35, 0.02, LINE)
    yy += 0.82
box(s, 9.7, 2.25, 2.65, 3.6, DARK)
text(s, 10.0, 2.6, 2.1, 2.9,
     [[("35–45", {"size": 34, "bold": True, "color": PAPER})],
      [("minutes, one CPU core", {"size": 12.5, "color": "8FA0B6"})],
      [(" ", {"size": 8})],
      [("0", {"size": 34, "bold": True, "color": GOLD})],
      [("GPUs. No Hugging Face. LoRA written from scratch in ~20 lines.",
        {"size": 12.5, "color": "8FA0B6"})]],
     space=1.2)
text(s, 0.85, 6.4, 11.5, 0.5,
     [[("Anyone in this room can reproduce all 48 runs on the laptop they "
        "brought.", {})]], size=15.5, bold=True, color=MUTED)
footer(s, "Setup")

# =========================================================== 12 VERIFY
s = slide(dark=True)
kicker(s, "Before measuring anything")
title(s, "Five checks before a single training step", size=36)
checks = ["gradients actually reach every LoRA adapter",
          "adapters sit on the real projections, not on copies",
          "base weights are genuinely frozen",
          "LoRA is a perfect no-op at initialisation",
          "the curriculum really does reorder the data"]
yy = 2.3
for i, c in enumerate(checks):
    chip(s, 0.85, yy, 0.42, 0.42, str(i + 1), GOLD, txtcol=DARK, size=14)
    text(s, 1.48, yy + 0.02, 6.6, 0.42, c, size=15.5, color="D6DFEA", anchor="m")
    yy += 0.6
box(s, 8.5, 2.3, 3.95, 4.3, DARKSF)
text(s, 8.85, 2.62, 3.25, 2.8,
     [[("Check 2 caught a silent aliasing bug.", {"bold": True, "color": GOLD,
                                                  "size": 17})],
      [(" ", {"size": 7})],
      [("The adapters existed, the parameter count looked right, the loss "
        "curves looked beautiful — and no gradient was flowing.",
        {"color": "B9C6D8", "size": 13.5})],
      [(" ", {"size": 7})],
      [("Those results were fiction. If I had not written verify.py I would "
        "be defending them today.", {"color": PAPER, "size": 13.5,
                                     "bold": True})]],
     space=1.25)
footer(s, "Setup")

# =========================================================== 13 RESULT 1
s = slide()
kicker(s, "Result 1 — all eight cells")
title(s, "The full method is not at the bottom", size=36)
pic(s, os.path.join(FIG, "fig2_factorial.png"), 1.75, 2.02, 9.84)
box(s, 0.85, 6.22, 11.5, 0.6, SOFT)
text(s, 1.2, 6.22, 11.0, 0.6,
     [[("Blue = naive baseline.  Red = full CG-PED.  ", {"color": BODY}),
       ("On real text the composed method is 0.254 nats worse than doing "
        "nothing clever at all.", {"bold": True, "color": RED})]],
     size=15, anchor="m")
footer(s, "Results")

# =========================================================== 14 RESULT 2
s = slide()
kicker(s, "Result 2 — the headline")
title(s, "They compose — almost perfectly", size=36)
pic(s, os.path.join(FIG, "fig3_effects.png"), 2.2, 1.96, 8.93)
yy = 5.62
statnum(s, 0.85, yy, 2.7, "−0.0003", "composition gap\nsynthetic", TEAL,
        numsize=27)
statnum(s, 3.75, yy, 2.7, "+0.0008", "composition gap\nreal text", TEAL,
        numsize=27)
statnum(s, 6.65, yy, 2.7, "±0.0055 / ±0.0099", "seed noise floor",
        MUTED, numsize=18)
statnum(s, 9.55, yy, 2.8, "12–18×", "smaller than the noise\nthat "
        "surrounds it", RED, numsize=27)
footer(s, "Results")

# =========================================================== 15 PRICE LIST
s = slide(soft=True)
kicker(s, "Result 3")
title(s, "Who helps, who costs, who is just noise", size=36)
hdr = [(0.85, 2.6, "TECHNIQUE"), (3.55, 1.5, "SYNTHETIC"), (5.25, 1.5, "REAL TEXT"),
       (7.05, 5.3, "VERDICT")]
for x, w, t in hdr:
    text(s, x, 2.16, w, 0.3, t, size=11.5, bold=True, color=MUTED)
box(s, 0.85, 2.5, 11.5, 0.025, INK)
rows = [
    ("Curriculum", AMBER, "−0.004", "+0.012",
     "Inside the noise, and it flips sign. Free — and worth what you pay."),
    ("LoRA", RED, "+0.027", "+0.278",
     "28× the noise on real text. The dominant term, and it is a cost."),
    ("Distillation", TEAL, "−0.007", "−0.037",
     "The only effect that helps on both benchmarks. The one thing to keep."),
]
yy = 2.65
for nm, c, sy, re_, verdict in rows:
    box(s, 0.85, yy, 11.5, 1.02, PAPER)
    box(s, 0.85, yy, 0.075, 1.02, c)
    text(s, 1.15, yy, 2.4, 1.02, nm, size=19, bold=True, color=c, anchor="m")
    text(s, 3.55, yy, 1.5, 1.02, sy, size=18, bold=True, color=INK, anchor="m",
         font=MONO)
    text(s, 5.25, yy, 1.6, 1.02, re_, size=18, bold=True,
         color=RED if re_.startswith("+") else TEAL, anchor="m", font=MONO)
    text(s, 7.05, yy, 5.15, 1.02, verdict, size=14, color=BODY, anchor="m",
         space=1.2)
    yy += 1.12
box(s, 0.85, 6.1, 11.5, 0.74, DARK)
text(s, 1.2, 6.1, 11.0, 0.74,
     [[("Sticky notes only save time if the book is already written. Training "
        "from scratch, the book is blank — and rank 4 has to write all of it.",
        {})]], size=15.5, bold=True, color=PAPER, anchor="m")
footer(s, "Results")

# =========================================================== 16 BUDGET
s = slide()
kicker(s, "Result 4 — the reality check")
title(s, "Equal steps flatters us. Equal wall clock does not.", size=34)
pic(s, os.path.join(FIG, "fig4_budget.png"), 2.35, 1.98, 8.64)
yy = 5.66
statnum(s, 0.85, yy, 2.7, "2.1–2.3×", "cost per step\nvs the naive run",
        RED, numsize=27)
statnum(s, 3.75, yy, 2.7, "131 / 139", "steps CG-PED affords\nin the naive "
        "run's time", RED, numsize=27)
statnum(s, 6.65, yy, 2.7, "3 / 3", "synthetic seeds catch up\n(by ~360–405 "
        "steps)", TEAL, numsize=27)
statnum(s, 9.55, yy, 2.8, "1 / 3", "real-text seeds catch up\nwithin 900 steps",
        RED, numsize=27)
footer(s, "Results")

# =========================================================== 17 SELF-ATTACK
s = slide(soft=True)
kicker(s, "We attacked our own result")
title(s, "“LoRA was never meant for training from scratch.”", size=32)
text(s, 0.85, 2.02, 11.5, 0.42,
     [[("Pretrain on Shakespeare → adapt to Jane Austen. Same character "
        "vocabulary, genuinely different register.", {})]],
     size=15.5, color=MUTED)
pic(s, os.path.join(SCR, "fig_transfer.png"), 1.95, 2.5, 9.42)
box(s, 0.85, 6.06, 11.5, 0.76, DARK)
text(s, 1.2, 6.06, 11.0, 0.76,
     [[("The objection was fair, and it did not save the method. LoRA's cost "
        "shrinks by 71% — to ", {"color": "B9C6D8"}),
       ("+0.082, still 3.7× the noise floor", {"bold": True, "color": GOLD}),
       (". Better is not free.", {"color": "B9C6D8"})]],
     size=15.5, anchor="m", space=1.2)
footer(s, "Results")

# =========================================================== 18 CONTRIBUTION
s = slide(dark=True)
kicker(s, "Contribution")
title(s, "Five things that did not exist before this study")
items = [
    ("The design", "the first full 2³ factorial over curriculum × LoRA × "
     "distillation — all eight cells exist"),
    ("The number", "the composition gap, defined and measured: "
     "−0.0003 / +0.0008 nats"),
    ("The price list", "each component's effect published against its own "
     "noise floor, so you can tell a result from seed luck"),
    ("The reporting", "equal-wall-clock alongside equal-step, because that is "
     "the budget practitioners actually have"),
    ("A clean negative", "with a mechanism — plus assert_lora_live, reusable "
     "by any paper that touches LoRA"),
]
yy = 2.18
for i, (h, b) in enumerate(items):
    chip(s, 0.85, yy, 0.44, 0.44, str(i + 1), GOLD, txtcol=DARK, size=15)
    text(s, 1.52, yy - 0.02, 10.6, 0.55,
         [[(h + "  —  ", {"bold": True, "color": PAPER}),
           (b, {"color": "A9B9CD"})]], size=15.5, space=1.2)
    yy += 0.8
box(s, 0.85, 6.08, 11.5, 0.7, DARKSF)
text(s, 1.2, 6.08, 11.0, 0.7,
     [[("The machinery is technique-agnostic. Swap in QLoRA, DoRA or data "
        "pruning and nothing but the flags changes.", {})]],
     size=15.5, bold=True, color=GOLD, anchor="m")
footer(s, "Contribution")

# =========================================================== 19 LIMITS
s = slide()
kicker(s, "Honest limitations")
title(s, "Where this result stops being true")
left = [
    "82K parameters, 300 steps, character-level — a laboratory, not production.",
    "Everything trained from scratch. LoRA is measured outside the regime it "
    "was designed for.",
    "Three seeds. Enough for a noise floor, not for a p-value.",
]
right = [
    "Two benchmarks, one of them synthetic and one of them 1.1 MB of "
    "Shakespeare.",
    "The curriculum's difficulty score is weak on real text — part of why "
    "that effect vanishes.",
    "Additivity is established here. Nobody has shown it holds at scale.",
]
bullets(s, 0.85, 2.35, 5.5, left, size=15.5, gap=1.05, dot=RED, wrap_h=1.0)
bullets(s, 6.85, 2.35, 5.5, right, size=15.5, gap=1.05, dot=RED, wrap_h=1.0)
box(s, 0.85, 5.75, 11.5, 0.95, SOFT)
text(s, 1.2, 5.75, 11.0, 0.95,
     [[("I would rather draw the boundary myself than have you find it. "
        "Every one of these is a next experiment, and the codebase already "
        "supports them.", {})]], size=15.5, color=INK, anchor="m", space=1.2)
footer(s, "Discussion")

# =========================================================== 20 CONCLUSION
s = slide(soft=True)
kicker(s, "Conclusion")
title(s, "What I'd tell a practitioner on Monday", size=36)
items = [
    (TEAL, "They add up. So stop running combinations.",
     "If effects are additive you can predict any of the eight cells from "
     "three single-factor runs. A combinatorial budget becomes a linear one."),
    (TEAL, "Keep distillation.",
     "The only component that helped on both benchmarks — if you can afford "
     "the teacher's compute."),
    (RED, "Use LoRA for memory, not for loss.",
     "And only on a pretrained backbone, where its price drops by 71% instead "
     "of dominating everything."),
    (AMBER, "Don't pay for a curriculum you can't measure.",
     "Ordering data helped nothing here, and the difficulty score is the "
     "prime suspect."),
]
yy = 2.2
for c, h, b in items:
    box(s, 0.85, yy, 11.5, 0.94, PAPER)
    box(s, 0.85, yy, 0.075, 0.94, c)
    text(s, 1.2, yy, 5.0, 0.94, h, size=16.5, bold=True, color=INK, anchor="m")
    text(s, 6.35, yy, 5.85, 0.94, b, size=13.5, color=BODY, anchor="m",
         space=1.2)
    yy += 1.0
box(s, 0.85, 6.32, 11.5, 0.02, INK)
text(s, 0.85, 6.45, 11.5, 0.5,
     [[("My own data says: do not use my composed method. That is the result, "
        "and I am reporting it as one.", {})]],
     size=16, bold=True, color=RED)
footer(s, "Conclusion")

# =========================================================== 21 FUTURE
s = slide(dark=True)
kicker(s, "What's next — and I want your verdict")
text(s, 0.85, 1.45, 7.6, 1.6,
     [[("TwinAI", {"color": GOLD})],
      [("How much of a life is enough?", {"color": PAPER, "size": 25})]],
     size=44, bold=True, space=1.1, space_after=6)
text(s, 0.85, 3.35, 7.6, 1.5,
     [[("Can an AI build a faithful behavioural model of ", {"color": "B9C6D8"}),
       ("one person", {"bold": True, "color": PAPER}),
       (" from their longitudinal history — and predict decisions they "
        "have never faced?", {"color": "B9C6D8"})]],
     size=18, space=1.3)
box(s, 0.85, 4.95, 7.6, 0.02, "3B4C68")
text(s, 0.85, 5.2, 7.6, 1.5,
     [[("The same question as this paper, moved from tokens to lifetimes: if "
        "10% of a history buys 70% of the fidelity, most of a life is "
        "redundant data. That is ", {"color": "8CA0BA"}),
       ("data-efficient AI, applied to people", {"bold": True, "color": GOLD}),
       (".", {"color": "8CA0BA"})]], size=15.5, space=1.28)
box(s, 8.85, 1.45, 3.6, 5.15, DARKSF)
text(s, 9.2, 1.8, 2.95, 0.35, "AFFORDABLE BY DESIGN", size=11.5, bold=True,
     color=GOLD)
bullets(s, 9.2, 2.35, 2.95, [
    "Phase 1 — synthetic individuals, known ground truth",
    "Phase 2 — small consented behavioural study",
    "No foundation model trained: retrieval + episodic memory + light "
    "fine-tuning",
    "Classical ML baselines first, LLM agents last",
], size=12.5, gap=0.95, color="B9C6D8", dot=GOLD, dotsize=0.085, wrap_h=0.9)
text(s, 9.2, 6.05, 2.95, 0.4, "Is it doable? That is my question to you.",
     size=12.5, bold=True, color=PAPER)
footer(s, "Future work")

# =========================================================== 22 CLOSE
s = slide(dark=True)
box(s, 0, 0, 0.28, 7.5, GOLD)
text(s, 1.0, 2.15, 10.6, 2.6,
     [[("Six months ago: Python and web development.", {"color": "8CA0BA"})],
      [("Today: 48 training runs and a negative result\nI can defend "
        "line by line.", {"color": PAPER})]],
     size=30, bold=True, space=1.18, space_after=16)
box(s, 1.0, 5.05, 1.4, 0.055, GOLD)
text(s, 1.0, 5.45, 10.6, 0.9,
     [[("Thank you.", {"size": 26, "bold": True, "color": GOLD})],
      [("Questions — and I would genuinely like your verdict on TwinAI.",
        {"size": 16, "color": "B9C6D8"})]], space=1.35, space_after=8)
footer(s, "CG-PED  ·  Abel Tezare")

# ================================================== SPEAKER NOTES
NOTES = [
 # 01
 "~20s. Don't read the title. Say: 'Three well-known tricks make training "
 "cheaper. Everybody uses them together. Nobody had checked whether their "
 "savings actually add up. I did — 48 training runs, one CPU core.' Then move.",
 # 02
 "~40s. Deliver this straight, not apologetically. The point is not modesty, "
 "it's provenance: this came out of the training programme. Land 'this deck "
 "is the receipt' and go. Flag that you will end with a proposal you want "
 "their opinion on — it buys attention for 20 minutes.",
 # 03
 "~30s. Pause after 'on a laptop.' This is the hook for the whole motivation. "
 "The one sentence that matters: nothing magic happened, a pile of efficiency "
 "techniques happened — and they are almost always reported one at a time.",
 # 04
 "~60s. Expect a question here: 'why not QLoRA / DoRA / something from this "
 "year?' Answer with the timeline: I needed techniques whose solo effects are "
 "already well-characterised, so that any surprise in the combination is "
 "attributable to the interaction, not to an unstable technique. The design "
 "is technique-agnostic — swapping in QLoRA changes flags, not machinery.",
 # 05
 "~70s. Spend the analogies here, they pay off for the rest of the talk. "
 "Curriculum = you don't start a maths course with calculus. LoRA = sticky "
 "notes in the margin instead of rewriting the book. Distillation = learning "
 "from a senior who already made the mistakes. Plant the LoRA one carefully — "
 "you cash it in on slide 15.",
 # 06
 "~45s. The checkout analogy is the whole motivation in one image. Then hit "
 "the right-hand column: CLPD, POCL, TSCL all compare a combination against "
 "an uncomposed baseline. That answers 'better than nothing?' — never 'better "
 "than its parts?'. This is the gap, in their own literature.",
 # 07
 "~60s. This is the intellectual core; slow down. Read the formula out loud. "
 "Combination = effects side by side, addition is honest. Composition = each "
 "technique changes the input of the next, so the sum is a guess. Concretely: "
 "the curriculum changes what LoRA sees; LoRA changes how much the student "
 "can absorb from the teacher.",
 # 08
 "~40s. Three worlds — synergy, additivity, interference. Four one-factor "
 "experiments produce the first three bars and then assume the fourth. You "
 "cannot distinguish these three pictures without running the combination. "
 "That is why the design has to be factorial.",
 # 09
 "~60s. Walk the eight cells left to right; name the naive baseline (blue "
 "outline) and the full method (dark, red outline). Then the arithmetic: "
 "8 cells x 3 seeds x 2 benchmarks = 48 runs. Stress the control: everything "
 "else identical — architecture, LR, batch, steps, seeds.",
 # 10
 "~45s. Anticipate 'four cells would have been cheaper.' Four gives you main "
 "effects and nothing else — no interaction term, and no independent estimate "
 "of noise. Eight costs exactly double and produces the composition gap plus "
 "the seed-spread noise floor. No cheaper design produces that number.",
 # 11
 "~45s. Keep this brisk — it is setup, not argument. Emphasise two things: "
 "the teacher is trained from scratch on the same data and then frozen (it is "
 "NOT a pretrained model), and the whole study runs in 35-45 minutes on one "
 "CPU core. Anyone in the room can reproduce it.",
 # 12
 "~50s. This slide buys you credibility for every number that follows. Then "
 "tell the bug story honestly: the adapters existed, the parameter count was "
 "right, the loss curves looked beautiful — and no gradient was flowing. "
 "Those results were fiction. verify.py is why I am not defending them today.",
 # 13
 "~45s. Point at blue (naive) and red (full CG-PED). On the synthetic "
 "benchmark everything is within a hair. On real text the composed method is "
 "0.254 nats WORSE than doing nothing clever. Do not soften it — the rest of "
 "the talk explains it.",
 # 14
 "~75s. The headline. Left panel: only one main effect is large, and it is "
 "the wrong sign. Right panel: observed minus additive prediction is "
 "-0.0003 and +0.0008, against seed noise of 0.0055 and 0.0099. The "
 "interaction is 12-18x smaller than the noise around it. They compose "
 "additively — which is exactly why the whole thing loses: nothing compounds, "
 "so you get the sum, and the sum is negative.",
 # 15
 "~70s. Read the table as a price list. Curriculum: inside the noise and it "
 "flips sign between benchmarks — free, and worth what you pay. LoRA: 28x the "
 "noise on real text, the dominant term, and it is a cost. Distillation: the "
 "only consistent win. Then cash in the analogy: sticky notes only save time "
 "if the book is already written. From scratch, the book is blank, and rank 4 "
 "has to write all of it.",
 # 16
 "~60s. Pre-empt the 'you flattered yourself with equal steps' objection by "
 "raising it first. A CG-PED step costs 2.1-2.3x a naive step because the "
 "teacher runs too. In the naive run's wall clock, CG-PED affords only "
 "131/139 steps. Synthetic catches up on all three seeds by ~400 steps; real "
 "text catches up on one seed out of three by 900.",
 # 17
 "~60s. Frame this as attacking your own result. The fair objection is that "
 "LoRA was designed for adapting a pretrained model, not for from-scratch "
 "training — so we pretrained on Shakespeare and adapted to Austen. LoRA's "
 "cost drops 71%, from +0.286 to +0.082. But 0.082 is still 3.7x the noise "
 "floor. Better is not free, and the recommendation does not flip.",
 # 18
 "~60s. If you only get one minute in the Q&A, this is the slide. The "
 "contribution is the design and the quantity it produces, not a new "
 "technique. Close with technique-agnosticism: QLoRA, DoRA, data pruning — "
 "the cells change, the machinery does not.",
 # 19
 "~45s. State the limits before the panel does. The two that matter: "
 "everything is from scratch, so LoRA is measured outside its home regime; "
 "and additivity is established here, at this scale — nobody has shown it "
 "holds at 7B. Each limitation maps to a next experiment the codebase "
 "already supports.",
 # 20
 "~60s. Land the practical takeaway. The additivity result is genuinely "
 "useful: if effects add, three single-factor runs predict all eight cells — "
 "a combinatorial budget becomes linear. Then the honest ending: my own data "
 "says do not use my composed method. Say it plainly; do not reframe it.",
 # 21
 "~75s. Shift energy — this is the ask, not a result. TwinAI is the same "
 "data-efficiency question moved from tokens to lifetimes: if 10% of a "
 "history buys 70% of the fidelity, most of a life is redundant data. Stress "
 "feasibility (synthetic individuals first, no foundation model trained) and "
 "ethics (consent, synthetic data first). Then ask them directly whether it "
 "is doable.",
 # 22
 "~20s. Close on the arc, then stop talking. Invite questions and explicitly "
 "invite their verdict on TwinAI — it turns Q&A into a conversation rather "
 "than an interrogation.",
]

assert len(NOTES) == len(prs.slides.__iter__.__self__._sldIdLst), "notes count"
for sl, nt in zip(prs.slides, NOTES):
    tf = sl.notes_slide.notes_text_frame
    tf.text = nt

prs.save(OUT)
print("saved", OUT, len(prs.slides.__iter__.__self__._sldIdLst), "slides")
